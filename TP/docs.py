from os import listdir
from config import CONFIG
import PyPDF2
import docx
import pdfplumber
from pptx import Presentation
import spacy
import re
from tika import parser
import mainFunctions
import probabilisticTopicRecognition


nlpSpanish = spacy.load("es_core_news_md")




def readPdf(filename):
    pdf = pdfplumber.open(filename)
    cant_pags = len(pdf.pages)
    string_lectura = ""
    for i in range(cant_pags):
        pagActual =  pdf.pages[i].extract_text()
        if pagActual is not None:
            string_lectura = string_lectura + pagActual

    pdf.close()
    return string_lectura

def readPdf2(filename):
    print(filename)
    archivo = open(filename, "rb")
    reader = PyPDF2.PdfFileReader(archivo)
    cant_pags = reader.getNumPages()
    print(cant_pags)
    string_lectura = ""
    print(string_lectura)
    for i in range(cant_pags):
        string_lectura = string_lectura + reader.getPage(i).extractText()

    return string_lectura


def readDocX(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)

def readDoc(filename):
    file = filename
    file_data = parser.from_file(file)
    text = file_data['content']

    return text

def readPpt(archivo):
    prs = Presentation(archivo)
    texto = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto = texto + shape.text + "/n"

    return texto


def readFile(path):
    fType = fileType(path)
    if fType == "pdf":
        return readPdf(path)

    elif fType == "docx":
        return readDocX(path)

    elif fType == "doc":
        return readDoc(path)

    elif fType == "pptx":
        return readPpt(path)

    else:
        return ""


class Parragraph:
    text = ""
    sentences = []
    metadata = []

    @staticmethod
    def newParagraph(text):
        par = Parragraph()
        par.text = text
        par.sentences = text.split(".")
        par.metadata = mainFunctions.preProcessor.preProcesar(text)
        return par


class Doc:
    name = ""
    text = ""
    fileType = ""
    author = ""
    paragraphs = []

    @staticmethod
    def newDoc(filePath):

        doc = Doc()
        doc.fileType = fileType(filePath)
        spl = filePath.split("/")
        doc.name = spl[len(spl)-1]
        doc.text = readFile(filePath)
        doc.paragraphs = [Parragraph.newParagraph(s) for s in doc.text.split("\n") if len(s) > 0]
        doc.paragraphs = [par for par in doc.paragraphs if len(par.metadata) > 0]
        return doc

    def generateParragraphMap(self):
        map = {}
        for parNum in range(len(self.paragraphs)):
            map[parNum] = []
        return map


    def generalPlagiarism(self, paragraphMap):
        total = 0

        count = len(paragraphMap)

        if CONFIG["QUESTIONS_SCORE"]:
            dividendo = count
        else:
            dividendo =len([ par for par in self.paragraphs if not isQuestion(par.text)])

        for i in range(count):
            if len(paragraphMap[i]) > 0:
                total += max(paragraphMap[i])


        return total/dividendo

    def getAuthors(self):
        people = find_persons(self.name + "/n" + self.text)
        if len(people) > 0:
            string = people[0]
            for p in people:
                if p != people[0]:
                    string = string + ", " + p
        else:
            string = "No se ha podido reconocer un autor."
        return string

    def getTopics(self):
        '''topics = find_topic(self.name + "/n" + self.text)
        if len(topics) > 0:
            string = topics[0]
            for t in topics:
                if t != topics[0]:
                    string = string + ", " + t
        else:
            string = "No se ha podida reconocer una temática."
            '''
        return probabilisticTopicRecognition.getTopicProbabilistic(self)



def ls(ruta):
    return listdir(ruta)


def fileType(path):
    pathTok = path.split(".")
    return (pathTok[len(pathTok) - 1]).lower()

def citations(string):
    text = string
    num_replaces = 100000000
    text = text.replace('“', '"', num_replaces).replace('”', '"', num_replaces).replace('„', '"', num_replaces).replace('‟', '"', num_replaces)
    pattern = r'\(([^"\)]*|\bAnónimo\b|"[^"\)]*")(, )([\d]+|n\.d\.|[\d]+[\w])\)|\[(1|2|3|4|5|6|7|8|9)(0|1|2|3|4|5|6|7|8|9)*\]'
    results = re.finditer(pattern, text)
    citations = [text[match.start(): match.end()] for match in results]
    return citations

def hasCitations(string):
    return len(citations(string)) > 0

def isQuestion(string):

    pattern = r'(0|1|2|3|4|5|6|7|8|9)+\)|(0|1|2|3|4|5|6|7|8|9)+-|(0|1|2|3|4|5|6|7|8|9)+\.'
    results = re.finditer(pattern, string)
    questionNum = [string[match.start(): match.end()] for match in results]

    return ("?" in string) or ("¿" in string) or (len(questionNum) > 0)


def obtainAuthor(string):
    propNouns = [x for x in nlpSpanish(string) if x.pos_=="PROPN"]
    nombre = propNouns[0].text + " " + propNouns[1].text
    return nombre


def find_persons(text):

    nlp = spacy.load("C:/respositorios_git/tp_nlp/TP/models/alumno")
    doc2 = nlp(text)

    persons = [ent.text for ent in doc2.ents if ent.label_ == 'ALUMNO']
    return persons

def find_topic(text):
    nlp = spacy.load("C:/respositorios_git/tp_nlp/TP/models/tematica")
    doc2 = nlp(text)
    topics = [ent.text for ent in doc2.ents if ent.label_ == 'TEMATICA']

    return topics


